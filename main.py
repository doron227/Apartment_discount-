import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pathlib import Path
import json
from requests import get
from dira_url_json import dira_url_json
import smtplib
import os
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication

url = "https://data.gov.il/api/3/action/datastore_search?resource_id=7c8255d0-49ef-49db-8904-4cf917586031&"
dira_json_data = dira_url_json(url)

def filter_dict_keys(d, keys):
    return {k: v for k, v in d.items() if k in keys}
file = Path("GUI_output.txt")

if file.exists():
    with open(file, "r") as file:
        data = json.load(file)
    settlements_list = data["settlments"]
    price_ranges = data['price ranges']
    float_price_ranges = [[float(price) for price in prices] for prices in price_ranges]
    data['price ranges'] = float_price_ranges
    df = pd.DataFrame(data)
    url_settlements = 'https://data.gov.il/api/3/action/datastore_search?resource_id=d4901968-dad3-4845-a9b0-a57d027f11ab&'
    settlements_data = dira_url_json(url_settlements)['result']['records']
    filtered_code = []
    filter_with_name = []
    for dictionary in settlements_data:
        for key, value in dictionary.items():
            if str(key).strip() in settlements_list or str(value).strip() in settlements_list:
                filtered_code.append(dictionary["סמל_ישוב"])
                filter_with_name.append(dictionary["סמל_ישוב"])
                filter_with_name.append(str(value).strip())
                break
    #print(filter_with_name)
    #aftwr getting LamasCode from dict
    dira_data = dira_json_data['result']['records']  # Gives a list of dictionaries. Each dictionary represent one lottery.
    lamas_code = filtered_code
    filtered_all_data = []
    for code in lamas_code:
        filtered_data = [d for d in dira_data if d.get('LamasCode') == str(code)]
        if not filtered_data:
            continue
        else:
            filtered_all_data.append(filtered_data)
    price_ranges = [[float(x) for x in inner] for inner in price_ranges]
    filtered_by_price = []
    for filtered_data in filtered_all_data:
        for data in filtered_data:
            code = data['LamasCode']
            price = float(data['PriceForMeter'].replace(',', ''))
            for i, price_range in enumerate(price_ranges):
                if price_range[0] <= price <= price_range[1]:
                    filtered_by_price.append(data)
                    break
    #print(filtered_by_price)

    desired_keys = ['LamasCode', 'Winners', 'Subscribers','SubscribersBenyMakom',
                    'LotteryNativeHousingUnits','LotteryHousingUnits','PriceForMeter','LamasName','LotteryStatusValue']
    filtered_dict = []
    for d in filtered_by_price:
        filtered_dict.append(filter_dict_keys(d, desired_keys))

    filter_by_status = []
    for dict in range(len(filtered_dict)):
        if filtered_dict[dict]['LotteryStatusValue'] == 'ההגרלה פורסמה':
            filter_by_status.append(filtered_dict[dict])
    #print(filter_by_status)


    def replace_dict_value(lst, filter_by_status):
        new_filtered_dict = []
        for dct in filter_by_status:
            temp_dict = dct.copy()
            for i in range(len(lst)):
                if str(lst[i]) in temp_dict.values():
                    key = [k for k, v in temp_dict.items() if str(v) == str(lst[i])][0]
                    if i < len(lst) - 1:
                        temp_dict[key] = lst[i + 1]
                        break
            new_filtered_dict.append(temp_dict)

        return new_filtered_dict

    new_filtered_dict = replace_dict_value(filter_with_name, filter_by_status)
    df = pd.DataFrame(new_filtered_dict)
    if df.empty:

        file = Path("GUI_output.txt")

        if file.exists():
            with open(file, "r") as file:
                data = json.load(file)
                email_id = data['email']

        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        mail_subject = "No lotto"
        mail_text = "No lotto in your city "
        message = 'Subject: {}\n\n{}'.format(mail_subject, mail_text)
        server.login('diradiscountwizard@gmail.com', 'ukvzimaqprnpgfhb')
        server.sendmail('&&&&&', email_id, message)
    else:
        df['PriceForMeter'] = df['PriceForMeter'].str.replace(',','').astype(float)
        df.plot(x='LamasCode', y='PriceForMeter', kind='scatter')
        plt.xlabel('The desired cities')
        plt.ylabel("PriceForMeter")
        plt.title('Price per square meter in each lottery and average')

        mean_PriceForMeter = df.groupby("LamasCode").mean()["PriceForMeter"]

        colors = plt.cm.viridis(np.linspace(0,1,len(mean_PriceForMeter)))

        for i, mean in enumerate(mean_PriceForMeter):
            plt.axhline(mean, color=colors[i], linestyle='--', label=mean_PriceForMeter.index[i])

        plt.legend()

        plt.savefig("prices.png")
        #########
        df['Winners'] = df['Winners'].str.replace(',', '').astype(float)
        df['Subscribers'] = df['Subscribers'].str.replace(',', '').astype(float)
        mean_win = df.groupby("LamasCode").mean()["Winners"]
        mean_sub = df.groupby("LamasCode").mean()["Subscribers"]
        mean_ratio = mean_win/mean_sub
        df['ratio'] = df['Winners']/df['Subscribers']
        df.plot(x='LamasCode', y='ratio', kind='scatter')
        plt.xlabel('The desired cities')
        plt.ylabel("winner to subscriber ratio")
        plt.title('winner to subscriber ratio in each lottery and average')

        colors = plt.cm.viridis(np.linspace(0, 1, len(mean_ratio)))



        for i, mean in enumerate(mean_ratio):
            plt.axhline(mean, color=colors[i], linestyle='--', label=mean_ratio.index[i])

        plt.legend()

        plt.savefig("win_ratio.png")


        prices = []
        for data in dira_data:
            prices.append(data['PriceForMeter'])
        def filter_list(lst):
            filtered_list = []
            for item in lst:
                if any(char.isdigit() or char == '.' for char in item):
                    filtered_item = ''.join(char for char in item if char.isdigit() or char == '.')
                    if filtered_item:
                        filtered_list.append(filtered_item)
            return filtered_list
        prices_filter = filter_list(prices)
        price_stat_all = []
        for p_data in prices_filter:
            if 4000 < float(p_data) < 50000:
                price_stat_all.append(float(p_data))
        mean_price = sum(price_stat_all)/len(price_stat_all)
        maximum = max(price_stat_all)
        minimum = min(price_stat_all)

        # power point Presentation that will have two plots and some changing statistical text
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        left = top = 0
        right = 0
        bottom = Inches(4)
        height0 = Inches(3.5)
        height2 = Inches(3.5)
        width = prs.slide_width / Inches(1)
        height = prs.slide_height / Inches(1)

        pic = slide.shapes.add_picture("prices.png", left, top, height=height0)
        pic2 = slide.shapes.add_picture("win_ratio.png", right, bottom, height=height2)

        bottom = Inches(2)
        # Add the text box to the slide
        textbox = slide.shapes.add_textbox(Inches(5), bottom, 1, 1)

        text_frame = textbox.text_frame
        data_mean_PriceForMeter = mean_PriceForMeter.to_frame()
        data_mean_PriceForMeter.index.name = None
        print(data_mean_PriceForMeter)


        text_frame.text = f"National average price per \n square meter in lotteries: {mean_price},\n 'The maximum" \
                          f" price was {maximum} \n The minimum was {minimum}." \
                          f"\nThe average in the cities you chose is:\n \n{data_mean_PriceForMeter}"

        prs.save("dira_presentation.pptx")
        file = Path("GUI_output.txt")

        if file.exists():
            with open(file, "r") as file:
                data = json.load(file)
                email_id = data['email']

            def mail_by_time():
                server = smtplib.SMTP('smtp.gmail.com', 587)
                server.starttls()
                mail_subject = f"Dira Lottery Wizard update"
                mail_text = "Your update is:"

                message = MIMEMultipart()
                message['Subject'] = mail_subject
                message.attach(MIMEText(mail_text))

                # Attach the presentation file
                current_directory = os.path.dirname(os.path.abspath(__file__))
                file_path = os.path.join(current_directory, 'dira_presentation.pptx')
                with open(file_path, 'rb') as f:
                    attachment = MIMEApplication(f.read(), _subtype='pptx')
                attachment.add_header('Content-Disposition', 'attachment', filename='dira_presentation.pptx')
                message.attach(attachment)
                server.login('diradiscountwizard@gmail.com', 'ukvzimaqprnpgfhb')
                server.sendmail('&&&&&', email_id, message.as_string())
                server.quit()
            # if there is match between the price requested by the user and the price of the lottry in the market ,
            # then send the email and it will work according to the time the user set with this,
            # the schedule file will activate this file else it will send an email that No lotto in your price range.
            if len(filtered_by_price) > 0:
                mail_by_time()
            else:
                server = smtplib.SMTP('smtp.gmail.com', 587)
                server.starttls()
                mail_subject = "No lotto"
                mail_text = "No lotto in your price range "
                message = 'Subject: {}\n\n{}'.format(mail_subject, mail_text)
                server.login('diradiscountwizard@gmail.com', 'ukvzimaqprnpgfhb')
                server.sendmail('&&&&&', email_id, message)
